import streamlit as st
import pandas as pd
import os
from pptx import Presentation
from pptx.util import Inches
from type import create_type_incident_plots
from vlqc import create_vlqc_plots
from varqc import create_varqc_plots
from cgqc import create_cgqc_plots
from csqc import create_csqc_plots
from cyqc import create_cyqc_plots
from obqc import create_obqc_plots
from piqc import create_piqc_plots
from anqc import create_anqc_plots
from fuqc import create_fuqc_plots

# Fonction pour générer le fichier PPT
def generate_ppt(image_buffers, output_file):
    prs = Presentation()
    
    for image_buffer in image_buffers:
        slide_layout = prs.slide_layouts[5]  # Mise en page vide
        slide = prs.slides.add_slide(slide_layout)
        left = Inches(1)
        top = Inches(1)
        pic = slide.shapes.add_picture(image_buffer, left, top, width=Inches(8), height=Inches(5))
    
    prs.save(output_file)

# Interface Streamlit
st.title("Analyse des Incidents - Génération de PPT et Visualisation des Graphiques")

# Uploader le fichier Excel
uploaded_file = st.file_uploader("Choisissez un fichier XLSX", type="xlsx")

if uploaded_file is not None:
    df = pd.read_excel(uploaded_file)

    # Créer des graphiques pour "Type d'incident"
    type_incident_buffers = create_type_incident_plots(df)
    # Créer des graphiques pour "Véhicule lent"
    vlqc_buffers = create_vlqc_plots(df)
    # Créer des graphiques pour "Véhicule arrêté"
    varqc_buffers = create_varqc_plots(df)
    # Créer des graphiques pour "Congestion"
    cgqc_buffers = create_cgqc_plots(df)
    # Créer des graphiques pour "Contresens"
    csqc_buffers = create_csqc_plots(df)
    # Créer des graphiques pour "Cycliste"
    cyqc_buffers = create_cyqc_plots(df)
    # Créer des graphiques pour "Objet"
    obqc_buffers = create_obqc_plots(df)
    # Créer des graphiques pour "Piéton"
    piqc_buffers = create_piqc_plots(df)
    # Créer des graphiques pour "Animal"
    anqc_buffers = create_anqc_plots(df)
    # Créer des graphiques pour "Fumée"
    fuqc_buffers = create_fuqc_plots(df)


    # Afficher les graphiques sur Streamlit
    st.image(type_incident_buffers[0], caption="Nombre d'incidents par type d'incident")
    st.image(type_incident_buffers[1], caption="Nombre d'incidents par type d'incident (sans Filtré Web et Inhibé Web)")
    st.image(vlqc_buffers[0], caption="Nombre d'incidents Véhicule lent (Partie 1)")
    st.image(vlqc_buffers[1], caption="Nombre d'incidents Véhicule lent (Partie 2)")
    st.image(varqc_buffers[0], caption="Nombre d'incidents Véhicule arrêté (Partie 1)")
    st.image(varqc_buffers[1], caption="Nombre d'incidents Véhicule arrêté (Partie 2)")
    st.image(cgqc_buffers[0], caption="Nombre d'incidents Congestion (Partie 1)")
    st.image(vlqc_buffers[1], caption="Nombre d'incidents Congestion (Partie 2)")
    st.image(csqc_buffers[0], caption="Nombre d'incidents Contresens (Partie 1)")
    st.image(varqc_buffers[1], caption="Nombre d'incidents Contresens (Partie 2)")   
    st.image(cyqc_buffers[0], caption="Nombre d'incidents Cycliste (Partie 1)")
    st.image(cyqc_buffers[1], caption="Nombre d'incidents Cycliste (Partie 2)")
    st.image(obqc_buffers[0], caption="Nombre d'incidents Objet (Partie 1)")
    st.image(obqc_buffers[1], caption="Nombre d'incidents Objet (Partie 2)")
    st.image(piqc_buffers[0], caption="Nombre d'incidents Piéton (Partie 1)")
    st.image(piqc_buffers[1], caption="Nombre d'incidents Piéton (Partie 2)")
    st.image(anqc_buffers[0], caption="Nombre d'incidents Animal (Partie 1)")
    st.image(anqc_buffers[1], caption="Nombre d'incidents Animal (Partie 2)")   
    st.image(fuqc_buffers[0], caption="Nombre d'incidents Fumée (Partie 1)")
    st.image(fuqc_buffers[1], caption="Nombre d'incidents Fumée (Partie 2)")   

    # Générer le fichier PPT
    output_dir = "output"
    os.makedirs(output_dir, exist_ok=True)
    ppt_file = os.path.join(output_dir, "report.pptx")

    # Combiner les buffers d'image pour le PPT
    all_buffers = type_incident_buffers + vlqc_buffers + varqc_buffers + cgqc_buffers + csqc_buffers + cyqc_buffers + obqc_buffers
    + piqc_buffers + anqc_buffers + fuqc_buffers 
    generate_ppt(all_buffers, ppt_file)

    st.success(f"Le fichier PPT a été généré : {ppt_file}")

    # Ajouter un lien pour télécharger le PPT
    with open(ppt_file, "rb") as file:
        st.download_button(
            label="Télécharger le rapport PPT",
            data=file,
            file_name="report.pptx",
            mime="application/vnd.openxmlformats-officedocument.presentationml.presentation"
        )
